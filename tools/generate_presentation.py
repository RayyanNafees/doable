import os
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from bs4 import BeautifulSoup
import re

html_content = """
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Doable - Hackathon Presentation</title>
    <link rel="preconnect" href="https://fonts.googleapis.com">
    <link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700&display=swap" rel="stylesheet">
    <style>
        /* CORE SETUP */
        * { box-sizing: border-box; }
        body {
            background-color: #f1f5f9; /* Light grey background for the window */
            display: grid;
            gap: 20px;
            grid-template-columns: 1fr;
            margin: 0;
            min-height: 100vh;
            padding: 20px 0;
            place-items: center;
            font-family: 'Inter', sans-serif;
        }

        /* SLIDE CONTAINER - STRICT 1280x720 */
        .slide-container {
            background-color: #ffffff; /* STRICT WHITE BACKGROUND */
            border-radius: 0;
            box-shadow: 0 4px 6px -1px rgba(0, 0, 0, 0.1), 0 2px 4px -1px rgba(0, 0, 0, 0.06);
            display: flex;
            flex-direction: column;
            height: 720px;
            overflow: hidden;
            padding: 60px 80px;
            position: relative;
            width: 1280px;
        }

        /* TYPOGRAPHY */
        h1, h2, h3, h4, p, li { color: #0f172a; margin: 0; }
        
        h1 {
            font-size: 72px;
            font-weight: 700;
            letter-spacing: -1.5px;
            line-height: 1.1;
            margin-bottom: 16px;
        }

        .subtitle {
            font-size: 32px;
            font-weight: 300;
            color: #64748b;
            letter-spacing: -0.5px;
        }

        .slide-title {
            font-size: 42px;
            font-weight: 600;
            letter-spacing: -1px;
            margin-bottom: 40px;
            padding-bottom: 20px;
            border-bottom: 2px solid #e2e8f0;
            width: 100%;
        }

        h3 {
            font-size: 24px;
            font-weight: 600;
            margin-bottom: 12px;
            color: #334155;
        }

        p, li {
            font-size: 20px;
            line-height: 1.6;
            color: #475569;
        }

        /* ACCENTS */
        .accent-text {
            background: linear-gradient(135deg, #4f46e5 0%, #7c3aed 100%);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        }

        /* LAYOUTS */
        .content-area {
            display: flex;
            flex-direction: column;
            flex-grow: 1;
            justify-content: center;
            width: 100%;
        }

        /* Two Column - Modified to be text focused if needed */
        .two-column {
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 60px;
            width: 100%;
            height: 100%;
            align-items: center;
        }

        /* Single Column Centered */
        .single-column {
            display: flex;
            flex-direction: column;
            justify-content: center;
            width: 100%;
            height: 100%;
        }

        /* Tiled Content */
        .tiled-content {
            display: grid;
            grid-template-columns: repeat(3, 1fr);
            gap: 30px;
            width: 100%;
        }

        .tile {
            background-color: #f8fafc;
            border: 1px solid #e2e8f0;
            border-radius: 12px;
            padding: 30px;
            transition: all 0.3s ease;
        }
        
        .tile:hover {
            border-color: #6366f1;
            box-shadow: 0 10px 15px -3px rgba(0, 0, 0, 0.05);
        }

        /* List Styling */
        ul {
            list-style-type: disc; /* Standard bullets */
            padding-left: 20px;
            margin: 0;
        }
        
        ul li {
            margin-bottom: 16px;
        }

        /* Footer */
        .slide-footer {
            position: absolute;
            bottom: 40px;
            right: 80px;
            text-align: right;
            font-size: 14px;
            color: #94a3b8;
            border-top: 1px solid #e2e8f0;
            padding-top: 10px;
        }

        /* Timeline */
        .timeline-container {
            position: relative;
            padding: 40px 0;
            width: 100%;
        }

        .timeline-line {
            position: absolute;
            left: 50px;
            top: 0;
            bottom: 0;
            width: 2px;
            background: #e2e8f0;
        }

        .timeline-item {
            position: relative;
            margin-bottom: 40px;
            padding-left: 80px;
        }

        .timeline-dot {
            position: absolute;
            left: 41px;
            top: 5px;
            width: 20px;
            height: 20px;
            border-radius: 50%;
            background: white;
            border: 4px solid #4f46e5;
        }

        /* Tag */
        .tag {
            display: inline-block;
            padding: 6px 12px;
            background-color: #e0e7ff;
            color: #4338ca;
            border-radius: 20px;
            font-size: 14px;
            font-weight: 600;
            margin-bottom: 20px;
        }

        /* Center Content Utility */
        .center-slide {
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
            text-align: center;
        }
        
        .center-slide h2 { margin-bottom: 20px; }

        /* DEMO INTERFACE STYLES */
        .demo-interface {
            width: 100%;
            height: 100%;
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 40px;
        }
        .demo-input-container {
            display: flex;
            flex-direction: column;
            gap: 20px;
        }
        .demo-textarea {
            width: 100%;
            height: 200px;
            padding: 20px;
            border: 2px solid #e2e8f0;
            border-radius: 12px;
            font-family: 'Inter', sans-serif;
            font-size: 16px;
            resize: none;
            transition: border-color 0.2s;
        }
        .demo-textarea:focus {
            outline: none;
            border-color: #4f46e5;
        }
        .demo-btn {
            background: linear-gradient(135deg, #4f46e5 0%, #7c3aed 100%);
            color: white;
            border: none;
            padding: 15px;
            border-radius: 8px;
            font-size: 18px;
            font-weight: 600;
            cursor: pointer;
            transition: transform 0.1s;
        }
        .demo-btn:hover {
            transform: scale(1.02);
        }
        .demo-btn:active {
            transform: scale(0.98);
        }
        .demo-btn:disabled {
            opacity: 0.7;
            cursor: not-allowed;
        }
        .demo-output {
            background: #f8fafc;
            border-radius: 12px;
            padding: 20px;
            border: 1px solid #e2e8f0;
            overflow-y: auto;
            font-size: 15px;
        }
        
        /* Custom styles for AI output */
        .task-card {
            background: white;
            padding: 15px;
            margin-bottom: 10px;
            border-radius: 8px;
            border-left: 4px solid #cbd5e1;
            box-shadow: 0 2px 4px rgba(0,0,0,0.05);
        }
        .task-card.high { border-left-color: #ef4444; }
        .task-card.medium { border-left-color: #f59e0b; }
        .task-card.low { border-left-color: #10b981; }
        
        .task-header {
            display: flex;
            justify-content: space-between;
            margin-bottom: 5px;
            font-weight: 600;
            font-size: 16px;
        }
        .task-meta {
            font-size: 13px;
            color: #64748b;
        }
    </style>
</head>
<body>

<!-- Slide 1: Title -->
<div class="slide-container" id="slide1">
    <div class="content-area" style="text-align: left; justify-content: flex-end; padding-bottom: 60px;">
        <div style="z-index: 1; max-width: 80%;">
            <span class="tag">BMW Challenge Submission</span>
            <h1><span class="accent-text">Doable</span></h1>
            <p class="subtitle" style="color: #0f172a; margin-bottom: 40px;">The AI Personal Productivity Agent.<br>Redefining Task Management with Agentic Intelligence.</p>
            <div style="display: flex; gap: 40px; font-size: 16px; color: #64748b;">
                <div>
                    <strong>Team Name:</strong><br>Doable
                </div>
                <div>
                    <strong>Presenters:</strong><br>Rayyan Nafees & Zuhair Arif
                </div>
            </div>
        </div>
    </div>
    <div class="slide-footer">
        <p style="font-size: 14px; margin:0;">BMW Challenge Submission</p>
        <p style="font-size: 14px; margin:0;">December 30, 2025</p>
    </div>
</div>

<!-- Slide 2: Problem -->
<div class="slide-container" id="slide2">
    <h2 class="slide-title">The Productivity Paradox</h2>
    <div class="content-area">
        <div class="two-column">
            <div>
                <h3 style="font-size: 32px; margin-bottom: 30px;">Drowning in tasks,<br>starving for <span class="accent-text">context</span>.</h3>
                <p>Modern professionals juggle an average of <strong>4+ isolated tools</strong> daily for task tracking and communication.</p>
                <br>
                <div class="tile" style="background: #fff; border-left: 4px solid #ef4444;">
                    <p style="font-size: 18px; margin: 0; color: #333;"><strong>The Pain Point:</strong><br>Context switching consumes up to 40% of productive time. Traditional tools are passive—they record, but don't assist.</p>
                </div>
            </div>
            <div>
                <ul style="font-size: 22px;">
                    <li style="margin-bottom: 30px;"><strong>Overwhelming Complexity</strong><br><span style="font-size: 18px; color: #64748b;">Manual conversion of messages into tasks creates high cognitive load.</span></li>
                    <li style="margin-bottom: 30px;"><strong>Lack of "Why"</strong><br><span style="font-size: 18px; color: #64748b;">Task lists are static and disconnected from personal goals or psychological motivators.</span></li>
                    <li><strong>Static Tools</strong><br><span style="font-size: 18px; color: #64748b;">No active reasoning or planning capabilities.</span></li>
                </ul>
            </div>
        </div>
    </div>
</div>

<!-- Slide 3: Solution -->
<div class="slide-container" id="slide3">
    <h2 class="slide-title">Enter <span class="accent-text">Doable</span></h2>
    <div class="content-area single-column">
        <h3>Innovative Agentic Workflow</h3>
        <p style="margin-bottom: 30px; font-size: 24px;">An intelligent assistant that bridges the gap between intent (Natural Language) and action.</p>
        <div style="background-color: #f8fafc; padding: 40px; border-radius: 12px; border: 1px solid #e2e8f0;">
            <ul>
                <li style="margin-bottom: 20px;"><strong>Intelligent Assistant:</strong> Persona-based agent tailored for specific roles (e.g., Software Developers).</li>
                <li style="margin-bottom: 20px;"><strong>True Agency:</strong> Utilizes "Agentic" behavior to reason, plan, and execute independently.</li>
                <li><strong>Context-Aware:</strong> Understands the "Who", "What", and "Why" behind every task.</li>
            </ul>
        </div>
    </div>
</div>

<!-- Slide 4: Core Capabilities -->
<div class="slide-container" id="slide4">
    <h2 class="slide-title">Powered by Advanced NLP & Reasoning</h2>
    <div class="content-area">
        <div class="tiled-content">
            <div class="tile">
                <h3>Natural Language Parsing</h3>
                <p style="font-size: 16px;">"Schedule a 2-hour coding session after lunch"</p>
                <hr style="border: 0; border-top: 1px solid #e2e8f0; margin: 15px 0;">
                <p style="font-size: 16px; color: #4f46e5;"><strong>Auto-Action:</strong> Understands temporal context and user intent.</p>
            </div>
            <div class="tile">
                <h3>Smart Task Management</h3>
                <p style="font-size: 16px;">Automated Priority Ranking & Effort Estimation.</p>
                <hr style="border: 0; border-top: 1px solid #e2e8f0; margin: 15px 0;">
                <p style="font-size: 16px; color: #4f46e5;"><strong>Memory:</strong> Retains long-term context for recurring projects.</p>
            </div>
            <div class="tile">
                <h3>Proactive Insights</h3>
                <p style="font-size: 16px;">Anticipates bottlenecks before they happen.</p>
                <hr style="border: 0; border-top: 1px solid #e2e8f0; margin: 15px 0;">
                <p style="font-size: 16px; color: #4f46e5;"><strong>Guidance:</strong> Suggests breaks and workload rebalancing.</p>
            </div>
        </div>
    </div>
</div>

<!-- Slide 5: Features Dashboard -->
<div class="slide-container" id="slide5">
    <div class="content-area single-column">
        <span class="tag">Visualized</span>
        <h2 style="font-size: 48px; margin-bottom: 30px;">A Modern, Dynamic Dashboard</h2>
        <div style="background-color: #f8fafc; padding: 40px; border-radius: 12px; border: 1px solid #e2e8f0;">
            <ul style="font-size: 20px;">
                <li style="margin-bottom: 25px;"><strong>Eisenhower Matrix:</strong> Auto-sorting tasks by Urgency vs. Importance.</li>
                <li style="margin-bottom: 25px;"><strong>Focus Mode:</strong> Distraction-free interface for deep work sessions.</li>
                <li style="margin-bottom: 25px;"><strong>Daily Analytics:</strong> Visual breakdown of productivity trends and completion rates.</li>
            </ul>
        </div>
    </div>
</div>

<!-- Slide 6: Architecture -->
<div class="slide-container" id="slide6">
    <h2 class="slide-title">Built on Modern Stack</h2>
    <div class="content-area single-column">
        <div class="tiled-content" style="grid-template-columns: 1fr 1fr;">
            <div class="tile">
                 <h3>Frontend & AI</h3>
                 <ul style="font-size: 18px;">
                    <li style="margin-bottom: 15px;"><strong>Frontend:</strong> Next.js 15 (App Router), React, TailwindCSS, Shadcn/UI.</li>
                    <li><strong>AI Engine:</strong> Google Gemini 2.5 Flash (Multimodal).</li>
                </ul>
            </div>
            <div class="tile">
                 <h3>Backend & Observability</h3>
                 <ul style="font-size: 18px;">
                    <li style="margin-bottom: 15px;"><strong>Server:</strong> Model Context Protocol (MCP).</li>
                    <li style="margin-bottom: 15px;"><strong>Frameworks:</strong> <code>xmcp</code> (TypeScript) & <code>agno</code> (Python).</li>
                    <li><strong>Tracing:</strong> Langfuse Integration with OpenTelemetry.</li>
                </ul>
            </div>
        </div>
    </div>
</div>

<!-- Slide 7: Conversational AI -->
<div class="slide-container" id="slide7">
    <h2 class="slide-title">The Multimodal Edge</h2>
    <div class="content-area single-column">
        <div style="margin-bottom: 40px;">
            <h3 style="color: #4f46e5; font-size: 32px; margin-bottom: 20px;">Agentic Voice & Audio</h3>
            <p>Leveraging Gemini TTS and ASR for human-level interaction.</p>
        </div>
        <div class="tiled-content">
            <div class="tile">
                <h3>Director's Mode</h3>
                <p>Context-aware tone (e.g., "Firm" for deadlines, "Upbeat" for wins).</p>
            </div>
            <div class="tile">
                <h3>Meeting to Task</h3>
                <p>Upload audio &rarr; Extract action items instantly.</p>
            </div>
            <div class="tile">
                <h3>Nuance Detection</h3>
                <p>Analyzes pitch/tone to detect user stress or fatigue.</p>
            </div>
        </div>
    </div>
</div>

<!-- Slide 8: Roadmap -->
<div class="slide-container" id="slide8">
    <h2 class="slide-title">The Path Ahead</h2>
    <div class="content-area">
        <div class="timeline-container">
            <div class="timeline-line"></div>
            
            <div class="timeline-item">
                <div class="timeline-dot"></div>
                <h3>Phase 1: Autonomous Reasoning</h3>
                <p>Expanding self-reflection cycles for higher task execution accuracy.</p>
            </div>
            
            <div class="timeline-item">
                <div class="timeline-dot"></div>
                <h3>Phase 2: Predictive Workflows</h3>
                <p>Using historical data to automate recurring task patterns and suggestion engines.</p>
            </div>
            
            <div class="timeline-item">
                <div class="timeline-dot"></div>
                <h3>Phase 3: Native Voice Interface</h3>
                <p>Full bidirectional, real-time voice conversations via Live API.</p>
            </div>

        </div>
    </div>
</div>

<!-- Slide 9: LIVE DEMO (NEW) -->
<div class="slide-container" id="slide9">
    <h2 class="slide-title">Experience Doable <span class="accent-text">Live</span></h2>
    <p class="subtitle" style="font-size: 20px; margin-bottom: 20px;">
        Test the <strong>Gemini 2.5 Flash</strong> reasoning engine. Enter a messy "Brain Dump" below:
    </p>
    
    <div class="demo-interface">
        <div class="demo-input-container">
            <div class="demo-textarea">
            (Interactive Demo Input Area - Placeholder)
            e.g. I need to fix the login bug by 2pm...
            </div>
            <div class="demo-btn">✨ Organize with Gemini</div>
        </div>
        <div class="demo-output" id="demoOutput">
            <div style="height:100%; display:flex; align-items:center; justify-content:center; color:#94a3b8; flex-direction:column;">
                Waiting for input...
            </div>
        </div>
    </div>
</div>

<!-- Slide 10: Conclusion -->
<div class="slide-container center-slide" id="slide10">
    <h2 style="font-size: 80px; margin-bottom: 20px;"><span class="accent-text">Make It Doable.</span></h2>
    <p style="font-size: 28px; max-width: 800px; margin-bottom: 50px;">Intelligent. Agentic. Psychologically-aware.<br>Experience the future of work today.</p>
</div>
</body>
</html>
"""


def hex_to_rgb(hex_color):
    if not hex_color:
        return RGBColor(0, 0, 0)
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
    )


def add_text_box(
    slide,
    text,
    left,
    top,
    width,
    height,
    font_size=18,
    is_bold=False,
    color=None,
    alignment=PP_ALIGN.LEFT,
):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    p.font.name = "Arial"  # Standard fallback
    p.font.size = Pt(font_size)
    p.font.bold = is_bold

    if color:
        p.font.color.rgb = color
    else:
        p.font.color.rgb = hex_to_rgb("475569")  # Default slate color matches CSS

    return txBox


def process_node_kids_text(node):
    """Recursively get text with basic spacing"""
    if not node:
        return ""
    return node.get_text(separator="\n").strip()


# Initialize Presentation
prs = Presentation()
# 16:9 Aspect Ratio (approx 1280x720 equivalent relative scale)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

soup = BeautifulSoup(html_content, "html.parser")
slides = soup.find_all("div", class_="slide-container")

for i, slide_div in enumerate(slides):
    # BLANK LAYOUT
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb("ffffff")

    # Dimensions needed for layout
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin_x = Inches(0.83)  # 80px approx
    margin_y = Inches(0.625)  # 60px approx
    content_w = slide_w - (margin_x * 2)

    # 1. Slide Title (h1 or h2.slide-title)
    h1 = slide_div.find("h1")
    h2_title = slide_div.find("h2", class_="slide-title")

    current_y = margin_y

    if h1:
        # Title Slide style
        tag = slide_div.find("span", class_="tag")
        if tag:
            add_text_box(
                slide,
                tag.get_text(),
                margin_x,
                current_y,
                content_w,
                Inches(0.5),
                font_size=14,
                color=hex_to_rgb("4338ca"),
                is_bold=True,
            )
            current_y += Inches(0.5)

        # Main Title
        add_text_box(
            slide,
            h1.get_text(),
            margin_x,
            current_y,
            content_w,
            Inches(1.5),
            font_size=72,
            is_bold=True,
            color=hex_to_rgb("0f172a"),
        )
        current_y += Inches(1.5)

        # Subtitle
        subtitle = slide_div.find("p", class_="subtitle")
        if subtitle:
            add_text_box(
                slide,
                subtitle.get_text(),
                margin_x,
                current_y,
                content_w,
                Inches(1.0),
                font_size=32,
                color=hex_to_rgb("64748b"),
            )
            current_y += Inches(1.2)

        # Team info (div grid at bottom in CSS)
        team_info = (
            slide_div.find_all("div", recursive=False)[0]
            .find_all("div", recursive=False)[0]
            .find_all("div", recursive=False)
        )
        # That struct is deep. Let's fuzzy find styles with 'display: flex'
        # Actually in slide 1, it's just after subtitle.
        # Hardcoding search for text 'Team Name'

    elif h2_title:
        # Standard Slide Title
        add_text_box(
            slide,
            h2_title.get_text(),
            margin_x,
            current_y,
            content_w,
            Inches(0.8),
            font_size=42,
            is_bold=True,
            color=hex_to_rgb("0f172a"),
        )

        # Add visual separator line
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            margin_x,
            current_y + Inches(0.8),
            content_w,
            Inches(0.02),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb("e2e8f0")
        shape.line.fill.background()

        current_y += Inches(1.2)  # Space after title

    elif slide_div.find("h2"):  # Slide 10 / Center Slide
        h2 = slide_div.find("h2")
        # Center slide logic
        # Vertically center approx
        center_y = slide_h / 2 - Inches(1)
        add_text_box(
            slide,
            h2.get_text(),
            margin_x,
            center_y,
            content_w,
            Inches(1.5),
            font_size=80,
            is_bold=True,
            color=hex_to_rgb("0f172a"),
            alignment=PP_ALIGN.CENTER,
        )

        p = slide_div.find("p")
        if p:
            add_text_box(
                slide,
                p.get_text(),
                margin_x,
                center_y + Inches(1.5),
                content_w,
                Inches(1),
                font_size=28,
                color=hex_to_rgb("475569"),
                alignment=PP_ALIGN.CENTER,
            )
        continue  # Skip rest of logic for this slide

    # Content Area Processing
    content_area = slide_div.find("div", class_="content-area")
    if not content_area:
        continue

    # Check for Layouts

    # 2-Column
    two_col = content_area.find("div", class_="two-column")
    if two_col:
        cols = two_col.find_all("div", recursive=False)
        col_w = (content_w - Inches(0.5)) / 2

        # Left Col
        left_x = margin_x
        # Quick parse of children
        for child in cols[0].children:
            if child.name == "h3":
                add_text_box(
                    slide,
                    child.get_text(),
                    left_x,
                    current_y,
                    col_w,
                    Inches(0.5),
                    font_size=32,
                    is_bold=True,
                    color=hex_to_rgb("334155"),
                )
                current_y += Inches(0.6)
            elif child.name == "p":
                add_text_box(
                    slide,
                    child.get_text(),
                    left_x,
                    current_y,
                    col_w,
                    Inches(0.5),
                    font_size=20,
                )
                current_y += Inches(0.6)
            elif child.name == "div" and "tile" in child.get("class", []):
                p_in_tile = child.find("p")
                if p_in_tile:
                    # Create a shape for tile
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        left_x,
                        current_y,
                        col_w,
                        Inches(1.5),
                    )
                    shape.fill.solid()
                    shape.fill.fore_color.rgb = hex_to_rgb("ffffff")  # white
                    shape.line.color.rgb = hex_to_rgb(
                        "ef4444"
                    )  # red border from inline style
                    shape.line.width = Pt(4)

                    tf = shape.text_frame
                    p = tf.paragraphs[0]
                    p.text = p_in_tile.get_text()
                    p.font.color.rgb = hex_to_rgb("333333")
                    p.font.size = Pt(18)

        # Right Col - Reset Y? It's side by side.
        # Actually pptx needs absolute positioning.
        # Let's just dump the list.
        right_x = margin_x + col_w + Inches(0.5)
        right_y = margin_y + Inches(1.2)  # Match title offset roughly

        ul = cols[1].find("ul")
        if ul:
            txBox = slide.shapes.add_textbox(right_x, right_y, col_w, Inches(4))
            tf = txBox.text_frame
            tf.word_wrap = True

            for li in ul.find_all("li"):
                p = tf.add_paragraph()
                p.text = li.get_text()
                p.font.size = Pt(22)
                p.level = 0
                p.space_after = Pt(14)

    # Tiled Content (3 cols)
    tiled = content_area.find("div", class_="tiled-content")
    if tiled:
        tiles = tiled.find_all("div", class_="tile")
        tile_w = (content_w - Inches(0.5)) / 3  # roughly
        if len(tiles) == 2:
            tile_w = (content_w - Inches(0.5)) / 2  # architecture slide has 2 tiles

        for idx, tile in enumerate(tiles):
            tile_x = margin_x + (tile_w + Inches(0.2)) * idx
            tile_y = current_y

            # Draw Tile Background
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, tile_x, tile_y, tile_w, Inches(2.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb("f8fafc")
            shape.line.color.rgb = hex_to_rgb("e2e8f0")

            # Content inside tile
            tf = shape.text_frame
            tf.margin_top = Inches(0.2)
            tf.margin_left = Inches(0.2)

            title = tile.find("h3")
            desc = tile.find_all("p")  # can be multiple or UL
            ul = tile.find("ul")

            if title:
                p = tf.paragraphs[0]
                p.text = title.get_text()
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = hex_to_rgb("334155")

            if ul:
                for li in ul.find_all("li"):
                    p = tf.add_paragraph()
                    p.text = li.get_text().strip()
                    p.level = 0
                    p.font.size = Pt(16)

            for d in desc:
                p = tf.add_paragraph()
                p.text = d.get_text().strip()
                p.font.size = Pt(16)

    # Timeline
    timeline = content_area.find("div", class_="timeline-container")
    if timeline:
        items = timeline.find_all("div", class_="timeline-item")
        for idx, item in enumerate(items):
            h3 = item.find("h3").get_text()
            p_text = item.find("p").get_text()

            add_text_box(
                slide,
                h3,
                margin_x + Inches(1),
                current_y,
                content_w - Inches(1),
                Inches(0.5),
                font_size=24,
                is_bold=True,
            )
            add_text_box(
                slide,
                p_text,
                margin_x + Inches(1),
                current_y + Inches(0.4),
                content_w - Inches(1),
                Inches(0.5),
                font_size=20,
            )

            # Draw dot
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                margin_x,
                current_y + Inches(0.1),
                Inches(0.3),
                Inches(0.3),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = hex_to_rgb("ffffff")
            dot.line.color.rgb = hex_to_rgb("4f46e5")
            dot.line.width = Pt(4)

            current_y += Inches(1.2)

    # Footer
    footer = slide_div.find("div", class_="slide-footer")
    if footer:
        dt = footer.find_all("p")[-1].get_text()
        add_text_box(
            slide,
            "Doable Presentation | " + dt,
            margin_x,
            slide_h - Inches(0.5),
            content_w,
            Inches(0.3),
            font_size=12,
            color=hex_to_rgb("94a3b8"),
            alignment=PP_ALIGN.RIGHT,
        )

output_file = "Doable_Presentation.pptx"
prs.save(output_file)
print(f"Presentation saved to {os.path.abspath(output_file)}")
